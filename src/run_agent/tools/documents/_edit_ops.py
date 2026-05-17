"""Targeted, in-place editing of generated documents.

`apply_edits` opens an existing file with the format-appropriate library and
applies an ordered list of operations, touching ONLY the parts each operation
names — everything else is preserved exactly. Returns the new bytes plus a
human-readable note for each operation (applied or skipped).

Operation vocabulary (validated per format):
  - replace_section {heading, content}
  - add_section     {heading, content, after?}
  - remove_section  {heading}
  - set_title       {title}
  - replace_text    {find, replace}
  - set_cell        {sheet?, cell, value}        (xlsx/csv)
  - add_row         {sheet?, values}             (xlsx/csv)
  - remove_row      {sheet?, row}                (xlsx/csv)
"""

import csv
import io
import re
from typing import Any

from run_agent.config import constants

_HEADING_LINE = re.compile(r"^(#{1,6})\s+(.*)$")
_CELL_REF = re.compile(r"^([A-Za-z]+)(\d+)$")


def apply_edits(
    file_type: str,
    content: bytes,
    operations: list[dict[str, Any]],
) -> tuple[bytes, list[str]]:
    """Dispatch to the handler for `file_type` and apply `operations`."""
    handler = _HANDLERS.get(file_type)
    if handler is None:
        raise ValueError(f"Editing is not supported for file type '{file_type}'")
    return handler(content, operations)


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------
def _norm(text: str | None) -> str:
    return (text or "").strip().lower()


def _content_paragraphs(content: str) -> list[str]:
    return [p.strip() for p in str(content).split("\n\n") if p.strip()]


def _parse_cell(ref: str) -> tuple[int, int]:
    """`B2` -> (column=2, row=2), both 1-based."""
    match = _CELL_REF.match(ref.strip())
    if not match:
        raise ValueError(f"Invalid cell reference '{ref}'")
    col = 0
    for char in match.group(1).upper():
        col = col * 26 + (ord(char) - 64)
    return col, int(match.group(2))


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------
def _edit_docx(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    from docx import Document

    doc = Document(io.BytesIO(content))
    notes: list[str] = []

    def is_heading(para: Any) -> bool:
        name = (para.style.name or "") if para.style else ""
        return name.startswith("Heading")

    def section_bounds(heading: str) -> tuple[int, int] | None:
        paras = doc.paragraphs
        for i, para in enumerate(paras):
            if is_heading(para) and _norm(para.text) == _norm(heading):
                end = len(paras)
                for j in range(i + 1, len(paras)):
                    if is_heading(paras[j]):
                        end = j
                        break
                return i, end
        return None

    for op in operations:
        kind = op.get("op")
        if kind == "set_title":
            target = next(
                (p for p in doc.paragraphs if p.style and p.style.name == "Title"),
                doc.paragraphs[0] if doc.paragraphs else None,
            )
            if target is not None:
                target.text = op.get("title", "")
                notes.append("set_title applied")
            else:
                notes.append("set_title skipped: no title paragraph")
        elif kind == "replace_section":
            bounds = section_bounds(op.get("heading", ""))
            if bounds is None:
                notes.append(f"replace_section skipped: '{op.get('heading')}' not found")
                continue
            i, end = bounds
            paras = doc.paragraphs
            anchor = paras[end] if end < len(paras) else None
            for body in paras[i + 1 : end]:
                body._element.getparent().remove(body._element)
            for text in _content_paragraphs(op.get("content", "")):
                if anchor is not None:
                    anchor.insert_paragraph_before(text)
                else:
                    doc.add_paragraph(text)
            notes.append(f"replace_section '{op.get('heading')}' applied")
        elif kind == "remove_section":
            bounds = section_bounds(op.get("heading", ""))
            if bounds is None:
                notes.append(f"remove_section skipped: '{op.get('heading')}' not found")
                continue
            i, end = bounds
            for para in doc.paragraphs[i:end]:
                para._element.getparent().remove(para._element)
            notes.append(f"remove_section '{op.get('heading')}' applied")
        elif kind == "add_section":
            after = op.get("after", "")
            anchor = None
            if after:
                bounds = section_bounds(after)
                if bounds is not None:
                    _, end = bounds
                    paras = doc.paragraphs
                    anchor = paras[end] if end < len(paras) else None
            if anchor is not None:
                anchor.insert_paragraph_before(op.get("heading", ""), style="Heading 1")
                for text in _content_paragraphs(op.get("content", "")):
                    anchor.insert_paragraph_before(text)
            else:
                doc.add_heading(op.get("heading", ""), level=1)
                for text in _content_paragraphs(op.get("content", "")):
                    doc.add_paragraph(text)
            notes.append(f"add_section '{op.get('heading')}' applied")
        elif kind == "replace_text":
            find, replace = op.get("find", ""), op.get("replace", "")
            hits = 0
            for para in doc.paragraphs:
                if find and find in para.text:
                    para.text = para.text.replace(find, replace)
                    hits += 1
            notes.append(f"replace_text applied to {hits} paragraph(s)")
        else:
            notes.append(f"unknown op '{kind}' skipped")

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue(), notes


# --------------------------------------------------------------------------
# PPTX  (a "section" is a slide, matched by its title)
# --------------------------------------------------------------------------
def _edit_pptx(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    notes: list[str] = []

    def slide_title(slide: Any) -> str:
        return slide.shapes.title.text if slide.shapes.title else ""

    def body_frame(slide: Any) -> Any:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0:
                return placeholder.text_frame
        return None

    def set_bullets(slide: Any, bullets: list[str]) -> bool:
        frame = body_frame(slide)
        if frame is None:
            return False
        frame.clear()
        for i, bullet in enumerate(bullets):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = bullet
        return True

    for op in operations:
        kind = op.get("op")
        if kind == "set_title":
            if prs.slides and prs.slides[0].shapes.title:
                prs.slides[0].shapes.title.text = op.get("title", "")
                notes.append("set_title applied")
            else:
                notes.append("set_title skipped: no title slide")
        elif kind in ("replace_section", "remove_section"):
            heading = op.get("heading", "")
            idx = next(
                (i for i, s in enumerate(prs.slides) if _norm(slide_title(s)) == _norm(heading)),
                None,
            )
            if idx is None:
                notes.append(f"{kind} skipped: slide '{heading}' not found")
                continue
            if kind == "remove_section":
                slide_list = prs.slides._sldIdLst
                slide_list.remove(list(slide_list)[idx])
                notes.append(f"remove_section '{heading}' applied")
            else:
                bullets = [b for b in str(op.get("content", "")).split("\n") if b.strip()]
                ok = set_bullets(prs.slides[idx], bullets)
                notes.append(
                    f"replace_section '{heading}' applied"
                    if ok
                    else f"replace_section skipped: slide '{heading}' has no body"
                )
        elif kind == "add_section":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title:
                slide.shapes.title.text = op.get("heading", "")
            set_bullets(slide, [b for b in str(op.get("content", "")).split("\n") if b.strip()])
            notes.append(f"add_section '{op.get('heading')}' applied")
        elif kind == "replace_text":
            find, replace = op.get("find", ""), op.get("replace", "")
            hits = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if find and find in run.text:
                                run.text = run.text.replace(find, replace)
                                hits += 1
            notes.append(f"replace_text applied to {hits} run(s)")
        else:
            notes.append(f"unknown op '{kind}' skipped")

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), notes


# --------------------------------------------------------------------------
# XLSX
# --------------------------------------------------------------------------
def _edit_xlsx(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content))
    notes: list[str] = []

    def sheet_for(op: dict[str, Any]) -> Any:
        name = op.get("sheet")
        if name and name in wb.sheetnames:
            return wb[name]
        return wb.active

    for op in operations:
        kind = op.get("op")
        try:
            if kind == "set_cell":
                ws = sheet_for(op)
                ws[op["cell"]] = op.get("value", "")
                notes.append(f"set_cell {op['cell']} on '{ws.title}' applied")
            elif kind == "add_row":
                ws = sheet_for(op)
                ws.append(list(op.get("values", [])))
                notes.append(f"add_row on '{ws.title}' applied")
            elif kind == "remove_row":
                ws = sheet_for(op)
                ws.delete_rows(int(op["row"]), 1)
                notes.append(f"remove_row {op['row']} on '{ws.title}' applied")
            elif kind == "replace_text":
                find, replace = op.get("find", ""), op.get("replace", "")
                hits = 0
                for ws in wb.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            if isinstance(cell.value, str) and find and find in cell.value:
                                cell.value = cell.value.replace(find, replace)
                                hits += 1
                notes.append(f"replace_text applied to {hits} cell(s)")
            else:
                notes.append(f"unsupported op '{kind}' for spreadsheet skipped")
        except Exception as exc:  # noqa: BLE001
            notes.append(f"{kind} skipped: {exc}")

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue(), notes


# --------------------------------------------------------------------------
# CSV
# --------------------------------------------------------------------------
def _edit_csv(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    text = content.decode("utf-8", errors="ignore")
    rows: list[list[str]] = list(csv.reader(io.StringIO(text)))
    notes: list[str] = []

    for op in operations:
        kind = op.get("op")
        try:
            if kind == "set_cell":
                col, row = _parse_cell(op["cell"])
                while len(rows) < row:
                    rows.append([])
                while len(rows[row - 1]) < col:
                    rows[row - 1].append("")
                rows[row - 1][col - 1] = str(op.get("value", ""))
                notes.append(f"set_cell {op['cell']} applied")
            elif kind == "add_row":
                rows.append([str(v) for v in op.get("values", [])])
                notes.append("add_row applied")
            elif kind == "remove_row":
                idx = int(op["row"]) - 1
                if 0 <= idx < len(rows):
                    del rows[idx]
                    notes.append(f"remove_row {op['row']} applied")
                else:
                    notes.append(f"remove_row skipped: row {op['row']} out of range")
            elif kind == "replace_text":
                find, replace = op.get("find", ""), op.get("replace", "")
                hits = 0
                for row in rows:
                    for i, val in enumerate(row):
                        if find and find in val:
                            row[i] = val.replace(find, replace)
                            hits += 1
                notes.append(f"replace_text applied to {hits} cell(s)")
            else:
                notes.append(f"unsupported op '{kind}' for CSV skipped")
        except Exception as exc:  # noqa: BLE001
            notes.append(f"{kind} skipped: {exc}")

    buffer = io.StringIO()
    csv.writer(buffer).writerows(rows)
    return buffer.getvalue().encode("utf-8"), notes


# --------------------------------------------------------------------------
# Markdown / plain text
# --------------------------------------------------------------------------
def _md_section_bounds(lines: list[str], heading: str) -> tuple[int, int] | None:
    """Find a Markdown section: its heading line through to the next heading."""
    for i, line in enumerate(lines):
        match = _HEADING_LINE.match(line)
        if match and _norm(match.group(2)) == _norm(heading):
            end = len(lines)
            for j in range(i + 1, len(lines)):
                if _HEADING_LINE.match(lines[j]):
                    end = j
                    break
            return i, end
    return None


def _edit_markdown(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    text = content.decode("utf-8", errors="ignore")
    notes: list[str] = []

    for op in operations:
        kind = op.get("op")
        if kind == "replace_text":
            find, replace = op.get("find", ""), op.get("replace", "")
            count = text.count(find) if find else 0
            text = text.replace(find, replace)
            notes.append(f"replace_text applied to {count} match(es)")
            continue

        lines = text.split("\n")
        if kind == "set_title":
            for i, line in enumerate(lines):
                match = _HEADING_LINE.match(line)
                if match and match.group(1) == "#":
                    lines[i] = f"# {op.get('title', '')}"
                    notes.append("set_title applied")
                    break
            else:
                lines.insert(0, f"# {op.get('title', '')}")
                notes.append("set_title applied (inserted)")
        elif kind == "replace_section":
            bounds = _md_section_bounds(lines, op.get("heading", ""))
            if bounds is None:
                notes.append(f"replace_section skipped: '{op.get('heading')}' not found")
            else:
                i, end = bounds
                lines[i + 1 : end] = ["", str(op.get("content", "")), ""]
                notes.append(f"replace_section '{op.get('heading')}' applied")
        elif kind == "remove_section":
            bounds = _md_section_bounds(lines, op.get("heading", ""))
            if bounds is None:
                notes.append(f"remove_section skipped: '{op.get('heading')}' not found")
            else:
                i, end = bounds
                del lines[i:end]
                notes.append(f"remove_section '{op.get('heading')}' applied")
        elif kind == "add_section":
            block = ["", f"## {op.get('heading', '')}", "", str(op.get('content', '')), ""]
            after = op.get("after", "")
            bounds = _md_section_bounds(lines, after) if after else None
            if bounds is not None:
                lines[bounds[1] : bounds[1]] = block
            else:
                lines.extend(block)
            notes.append(f"add_section '{op.get('heading')}' applied")
        else:
            notes.append(f"unknown op '{kind}' skipped")
        text = "\n".join(lines)

    return text.encode("utf-8"), notes


def _edit_txt(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    text = content.decode("utf-8", errors="ignore")
    notes: list[str] = []
    for op in operations:
        if op.get("op") == "replace_text":
            find, replace = op.get("find", ""), op.get("replace", "")
            count = text.count(find) if find else 0
            text = text.replace(find, replace)
            notes.append(f"replace_text applied to {count} match(es)")
        else:
            notes.append(
                f"op '{op.get('op')}' skipped: plain text supports only replace_text"
            )
    return text.encode("utf-8"), notes


# --------------------------------------------------------------------------
# PDF  (cannot be edited in place — text is extracted, edited, re-rendered)
# --------------------------------------------------------------------------
def _edit_pdf(content: bytes, operations: list[dict[str, Any]]) -> tuple[bytes, list[str]]:
    from pypdf import PdfReader

    from run_agent.tools.documents._render import render_pdf

    reader = PdfReader(io.BytesIO(content))
    text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
    notes: list[str] = [
        "note: PDF edits reflow the document — original layout is not preserved"
    ]

    for op in operations:
        if op.get("op") == "replace_text":
            find, replace = op.get("find", ""), op.get("replace", "")
            count = text.count(find) if find else 0
            text = text.replace(find, replace)
            notes.append(f"replace_text applied to {count} match(es)")
        else:
            notes.append(
                f"op '{op.get('op')}' skipped: PDF editing supports only replace_text "
                f"(extract the content and recreate the PDF for structural changes)"
            )

    sections = [{"heading": "", "content": text}]
    return render_pdf("", sections), notes


_HANDLERS = {
    constants.MIME_DOCX: _edit_docx,
    constants.MIME_PPTX: _edit_pptx,
    constants.MIME_XLSX: _edit_xlsx,
    constants.MIME_CSV: _edit_csv,
    constants.MIME_MD: _edit_markdown,
    constants.MIME_TXT: _edit_txt,
    constants.MIME_PDF: _edit_pdf,
}
