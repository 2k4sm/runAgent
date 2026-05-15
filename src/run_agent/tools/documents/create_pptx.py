"""PPTX creation tool."""

import io
from typing import Any

from pptx import Presentation

from run_agent.config.constants import MIME_PPTX
from run_agent.tools.base import BaseTool
from run_agent.tools.documents._common import filename_with_ext, finalize


class CreatePptxTool(BaseTool):
    name = "create_pptx"
    description = "Create a PowerPoint presentation (.pptx) with a title and slides."

    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Presentation title"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                        "required": ["heading"],
                    },
                    "description": "Content slides, each with a heading and bullet points",
                },
                "filename": {
                    "type": "string",
                    "description": "Output filename (without extension)",
                    "default": "presentation",
                },
            },
            "required": ["title", "slides"],
        }

    async def execute(
        self,
        title: str,
        slides: list[dict],
        filename: str = "presentation",
        _context: dict | None = None,
        **_kwargs: Any,
    ) -> str:
        presentation = Presentation()

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title

        for slide in slides:
            layout = presentation.slide_layouts[1]
            pptx_slide = presentation.slides.add_slide(layout)
            pptx_slide.shapes.title.text = slide["heading"]
            body = pptx_slide.placeholders[1].text_frame
            bullets = slide.get("bullets", [])
            for i, bullet in enumerate(bullets):
                paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
                paragraph.text = bullet

        buffer = io.BytesIO()
        presentation.save(buffer)
        return await finalize(
            filename_with_ext(filename, "pptx"),
            buffer.getvalue(),
            MIME_PPTX,
            _context,
        )
