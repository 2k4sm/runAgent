"""Read back an existing document's content and structure."""

import csv
import io
import json
from typing import Any

from run_agent.config import constants
from run_agent.services.file_service import FileService
from run_agent.tools.base import BaseTool

_TEXT_CAP = 8000


def _extract(file_type: str, content: bytes) -> dict[str, Any]:
    """Return a best-effort structured view of the document."""
    if file_type == constants.MIME_DOCX:
        from docx import Document

        doc = Document(io.BytesIO(content))
        sections: list[dict[str, Any]] = []
        for para in doc.paragraphs:
            style = (para.style.name or "") if para.style else ""
            if style.startswith("Heading") or style == "Title":
                sections.append({"heading": para.text, "style": style, "body": ""})
            elif para.text.strip():
                if not sections:
                    sections.append({"heading": "", "style": "", "body": ""})
                sections[-1]["body"] = (
                    f"{sections[-1]['body']}\n{para.text}".strip()
                )
        return {"sections": sections}

    if file_type == constants.MIME_PPTX:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides: list[dict[str, Any]] = []
        for slide in prs.slides:
            title = slide.shapes.title.text if slide.shapes.title else ""
            bullets: list[str] = []
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                    bullets += [
                        p.text for p in ph.text_frame.paragraphs if p.text.strip()
                    ]
            slides.append({"title": title, "bullets": bullets})
        return {"slides": slides}

    if file_type == constants.MIME_XLSX:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True)
        sheets: list[dict[str, Any]] = []
        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            sheets.append({
                "name": ws.title,
                "headers": list(rows[0]) if rows else [],
                "row_count": len(rows),
                "sample_rows": [list(r) for r in rows[1:6]],
            })
        return {"sheets": sheets}

    if file_type == constants.MIME_PDF:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return {"pages": [p[:_TEXT_CAP] for p in pages]}

    if file_type == constants.MIME_CSV:
        rows = list(csv.reader(io.StringIO(content.decode("utf-8", errors="ignore"))))
        return {
            "headers": rows[0] if rows else [],
            "row_count": len(rows),
            "sample_rows": rows[1:6],
        }

    # Markdown / plain text and anything else: return raw text.
    return {"text": content.decode("utf-8", errors="ignore")[:_TEXT_CAP]}


class ReadDocumentTool(BaseTool):
    name = "read_document"
    description = (
        "Read an existing document's current content and structure (section "
        "headings, slide titles, sheet layout). Call this before editing so your "
        "edit operations target the right sections."
    )

    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "asset_id": {
                    "type": "string",
                    "description": "The asset_id of the document (from list_documents).",
                },
            },
            "required": ["asset_id"],
        }

    async def execute(
        self,
        asset_id: str,
        _context: dict | None = None,
        **_kwargs: Any,
    ) -> str:
        if not _context or not _context.get("user_id"):
            return json.dumps({"status": "error", "message": "No run context."})
        try:
            content, asset = await FileService().download_asset(
                asset_id, _context["user_id"]
            )
        except ValueError as exc:
            return json.dumps({"status": "error", "message": str(exc)})

        return json.dumps({
            "status": "success",
            "asset_id": asset_id,
            "file_name": asset["file_name"],
            "file_type": asset["file_type"],
            "structure": _extract(asset["file_type"], content),
        })
